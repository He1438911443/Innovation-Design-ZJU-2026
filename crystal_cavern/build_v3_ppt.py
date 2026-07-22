#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""重建答辩 PPT v3 — 12 页深色科技风，嵌入真实图片（非占位符）。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image
import copy

# ---------- palette ----------
BG_TOP    = RGBColor(0x0F, 0x0C, 0x29)   # #0f0c29
BG_BOT    = RGBColor(0x30, 0x2B, 0x63)   # #302b63
PANEL     = RGBColor(0x1B, 0x17, 0x3A)   # deeper panel
PANEL2    = RGBColor(0x24, 0x1F, 0x4E)
CYAN      = RGBColor(0x22, 0xE1, 0xD4)   # cyan accent
CYAN_DIM  = RGBColor(0x12, 0xA8, 0xA0)
PURPLE    = RGBColor(0xB0, 0x6A, 0xFF)   # purple accent
PURPLE_DIM= RGBColor(0x7A, 0x4A, 0xC4)
GOLD      = RGBColor(0xFF, 0xC8, 0x4D)
WHITE     = RGBColor(0xF5, 0xF6, 0xFA)
GREY      = RGBColor(0xB8, 0xBC, 0xCC)
GREY_DIM  = RGBColor(0x80, 0x84, 0x9E)

FONT_CJK = "微软雅黑"
FONT_EN  = "Arial"
FONT_MONO= "Consolas"

EMU_W = 12191695
EMU_H = 6858000

prs = Presentation()
prs.slide_width  = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

ASSETS = "_ppt_assets"
REFDIR = "../水晶实景参考"

# ---------- helpers ----------
def add_bg(slide):
    """深色对角渐变背景：全屏矩形 + 手写 gradFill XML。"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    shp.line.fill.background()
    shp.shadow.inherit = False
    spPr = shp._element.spPr
    # remove existing fill
    for tag in ('a:noFill','a:solidFill','a:gradFill','a:blipFill','a:pattFill','a:grpFill'):
        e = spPr.find(qn(tag))
        if e is not None:
            spPr.remove(e)
    grad = spPr.makeelement(qn('a:gradFill'), {})
    gsLst = grad.makeelement(qn('a:gsLst'), {})
    def gs(pos, hexc):
        node = grad.makeelement(qn('a:gs'), {'pos': str(pos)})
        c = grad.makeelement(qn('a:srgbClr'), {'val': hexc})
        node.append(c)
        return node
    gsLst.append(gs(0,   '0F0C29'))
    gsLst.append(gs(50000, '1F1B47'))
    gsLst.append(gs(100000, '302B63'))
    grad.append(gsLst)
    lin = grad.makeelement(qn('a:lin'), {'ang':'2700000','scaled':'1'})
    grad.append(lin)
    # insert after geometry, before line
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        ln.addprevious(grad)
    else:
        spPr.append(grad)
    return shp

def _set_font(run, size, color, bold=False, font=FONT_CJK, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    # CJK font hint
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin','a:cs','a:ea'):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
    for tag in ('a:latin','a:cs','a:ea'):
        e = rPr.find(qn(tag))
        e.set('typeface', font)

def add_text(slide, l, t, w, h, text, size=18, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_CJK,
             line_spacing=1.15, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        _set_font(r, size, color, bold=bold, font=font, italic=italic)
    return tb

def add_round(slide, l, t, w, h, fill=PANEL, line=CYAN, line_w=1.0,
              radius=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    # adjust radius
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def add_rect(slide, l, t, w, h, fill=PANEL, line=None, line_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def add_line(slide, l, t, w, color=CYAN, weight=2.0):
    shp = slide.shapes.add_connector(1, l, t, l + w, t)
    shp.line.color.rgb = color
    shp.line.width = Pt(weight)
    return shp

def add_title(slide, text, color=WHITE):
    """Standard slide top: kicker bar + big title + cyan underline."""
    # accent bar
    bar = add_rect(slide, Inches(0.6), Inches(0.45), Inches(0.12), Inches(0.5),
                   fill=CYAN)
    add_text(slide, Inches(0.85), Inches(0.40), Inches(11.5), Inches(0.62),
             text, size=30, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # underline
    add_line(slide, Inches(0.62), Inches(1.12), Inches(12.1), CYAN, 2.0)
    # page corner gloss
    add_text(slide, Inches(11.6), Inches(0.42), Inches(1.4), Inches(0.5),
             "", size=11, color=GREY_DIM, align=PP_ALIGN.RIGHT)  # placeholder

def add_pagenum(slide, n):
    add_text(slide, Inches(12.0), Inches(7.05), Inches(1.1), Inches(0.35),
             f"{n:02d} / 12", size=10, color=GREY_DIM, align=PP_ALIGN.RIGHT,
             font=FONT_EN)
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.35),
             "Crystal Cavern  ·  ZJU 2026", size=10, color=GREY_DIM,
             font=FONT_EN)

def add_image_fit(slide, path, l, t, w, h):
    """Insert image fit into box (l,t,w,h) preserving aspect, centered, on dark frame."""
    img = Image.open(path)
    iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        nw = w
        nh = int(w / img_ratio)
    else:
        nh = h
        nw = int(h * img_ratio)
    nl = l + (w - nw) // 2
    nt = t + (h - nh) // 2
    # frame behind
    frm = add_rect(slide, nl - Emu(20000), nt - Emu(20000), nw + Emu(40000),
                   nh + Emu(40000), fill=PANEL, line=CYAN_DIM, line_w=1.0)
    pic = slide.shapes.add_picture(path, nl, nt, nw, nh)
    return pic, (nl, nt, nw, nh)

def add_caption(slide, l, t, w, text, color=GREY, size=11, align=PP_ALIGN.CENTER):
    add_text(slide, l, t, w, Inches(0.3), text, size=size, color=color,
             align=align, italic=True)

# =====================================================================
def slide(i):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    return s

# ---- SLIDE 1: TITLE ----
s = slide(1)
# decorative circle (semi-transparent purple glow)
c1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(-2.0), Inches(6.4), Inches(6.4))
c1.line.fill.background()
c1.shadow.inherit = False
_sp = c1._element.spPr
for tag in ('a:solidFill',):
    e = _sp.find(qn(tag))
    if e is not None:
        _sp.remove(e)
sf = _sp.makeelement(qn('a:solidFill'), {})
clr = sf.makeelement(qn('a:srgbClr'), {'val': '3A2F6B'})
alpha = clr.makeelement(qn('a:alpha'), {'val': '38000'})
clr.append(alpha); sf.append(clr)
ln = _sp.find(qn('a:ln'))
if ln is not None:
    ln.addprevious(sf)
else:
    _sp.append(sf)
# second smaller cyan ring
c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(-1.0), Inches(3.6), Inches(3.6))
c2.line.color.rgb = CYAN_DIM; c2.line.width = Pt(1.2)
c2.fill.background(); c2.shadow.inherit = False
# kicker
add_text(s, Inches(0.9), Inches(1.2), Inches(8), Inches(0.5),
         "◇  PROCEDURAL · MAYA PYTHON · 2026", size=14, color=CYAN,
         font=FONT_EN, bold=True)
add_text(s, Inches(0.85), Inches(1.75), Inches(11.5), Inches(1.6),
         "CRYSTAL CAVERN", size=66, color=WHITE, bold=True, font=FONT_EN)
add_text(s, Inches(0.9), Inches(3.25), Inches(11), Inches(0.7),
         "Python 驱动 Maya 程序化水晶洞穴生成系统", size=24, color=CYAN, bold=True)
# divider
add_line(s, Inches(0.95), Inches(4.05), Inches(6.5), PURPLE, 2.5)
add_text(s, Inches(0.92), Inches(4.2), Inches(11), Inches(0.45),
         "未来创新设计  ·  浙江大学 2026  ·  Prof. Xiaosong Yang",
         size=15, color=GREY)
add_text(s, Inches(0.92), Inches(4.65), Inches(11), Inches(0.45),
         "一键生成  ·  全自动管线  ·  零手工建模  ·  算法驱动",
         size=14, color=GREY_DIM)
# stat chips
chips = [("7", "Python 模块"), ("4", "经典算法"), ("56", "晶体/场景"),
         ("960", "帧飞越动画")]
cx = Inches(0.92); cy = Inches(5.55)
cw = Inches(2.7); ch = Inches(0.95); gap = Inches(0.18)
for i,(num,lab) in enumerate(chips):
    x = cx + i*(cw+gap)
    add_round(s, x, cy, cw, ch, fill=PANEL, line=PURPLE_DIM, line_w=1.0, radius=0.15)
    add_text(s, x, cy+Inches(0.08), cw, Inches(0.5), num, size=26, color=CYAN,
             bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(s, x, cy+Inches(0.55), cw, Inches(0.35), lab, size=11, color=GREY,
             align=PP_ALIGN.CENTER)
add_pagenum(s, 1)

# ---- SLIDE 2: 概念概览 (concepts + rendered result) ----
s = slide(2)
add_title(s, "从实景到程序化生成")
# left concept image (arnold render)
add_text(s, Inches(0.6), Inches(1.4), Inches(6), Inches(0.4),
         "◆ 算法生成结果  ·  最终渲染", size=13, color=CYAN, bold=True)
add_image_fit(s, "final_output/arnold_render.png",
              Inches(0.6), Inches(1.85), Inches(6.5), Inches(4.5))
add_caption(s, Inches(0.6), Inches(6.4), Inches(6.5),
            "紫水晶洞穴 (seed=2026) · 56 crystals · 1920×1080 Arnold 渲染")
# right text column
rx = Inches(7.45)
add_text(s, rx, Inches(1.45), Inches(5.3), Inches(0.5),
         "一句话目标", size=15, color=PURPLE, bold=True)
add_text(s, rx, Inches(1.95), Inches(5.3), Inches(1.1),
         "一键点击 → Maya 自动生成完整水晶洞穴：地形、晶体、灯光、飞越相机，全程零手工建模。",
         size=14, color=WHITE, line_spacing=1.3)
bullets = [
    "◉  全程序化 — 随机 seed 每次生成不同洞穴",
    "◉  参数驱动 — 密度 / 粗糙度 / 颜色 seed 全可控",
    "◉  算法引用 — 4 个已发表论文，规范标注",
    "◉  视觉目标 — SSS 辉光 + 体积雾 + 三层布光",
    "◉  可复现 — 同 seed = 相同结果，改参 = 新洞穴",
]
by = Inches(3.15)
for b in bullets:
    add_round(s, rx, by, Inches(5.3), Inches(0.55), fill=PANEL,
              line=PANEL2, line_w=0.75, radius=0.18)
    add_text(s, rx+Inches(0.25), by, Inches(5.0), Inches(0.55), b,
             size=12.5, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    by += Inches(0.66)
add_pagenum(s, 2)

# ---- SLIDE 3: 实景参考 (real cave refs) ----
s = slide(3)
add_title(s, "实景参考 — 墨西哥奈卡水晶洞")
refs = [
    "3b80ff12a708b83f91132dc8a65a690e~tplv-a9rns2rl98-pc_smart_face_crop-v1_512_384.png",
    "266362bbf6432235b7c4c433d7b25fd1~tplv-a9rns2rl98-pc_smart_face_crop-v1_512_384.jpeg",
    "2c4bcdb667d84c0e925a7a655bc23aa5~tplv-a9rns2rl98-pc_smart_face_crop-v1_512_384.jpeg",
    "a6a3556f6a16cb62ef7a7b628e70c3dc~tplv-a9rns2rl98-pc_smart_face_crop-v1_512_384.jpeg",
]
# 4 ref thumbnails row
tw = Inches(2.85); th = Inches(2.15); gap = Inches(0.13)
sx = Inches(0.6); sy = Inches(1.45)
for i,f in enumerate(refs):
    p = f"{REFDIR}/{f}"
    x = sx + i*(tw+gap)
    add_image_fit(s, p, x, sy, tw, th)
    add_text(s, x, sy+th+Inches(0.04), tw, Inches(0.3),
             f"参考 {i+1}", size=10, color=GREY, align=PP_ALIGN.CENTER, italic=True)
# arrow down to result
add_text(s, Inches(0.6), Inches(3.95), Inches(12), Inches(0.45),
         "▼  从实景到算法生成目标  ▼", size=15, color=CYAN, bold=True,
         align=PP_ALIGN.CENTER)
# target render (concept gen image) captioned
add_text(s, Inches(0.6), Inches(4.45), Inches(12), Inches(0.4),
         "◆ 算法目标 — 体积光 + 多色晶簇 + 次表面辉光", size=13, color=PURPLE,
         bold=True, align=PP_ALIGN.CENTER)
add_image_fit(s, "../生成图/0101.jpg", Inches(3.4), Inches(4.9),
              Inches(6.5), Inches(2.1))
add_pagenum(s, 3)

# ---- SLIDE 4: 算法创新 四卡 ----
s = slide(4)
add_title(s, "算法创新 — 四个经典算法的三维应用")
cards = [
    ("Diamond-Square", "Fournier, Fussell & Carpenter · 1982",
     "66,049顶点高分辨率地形\n洞穴塑造：径向幂+正弦扰动", CYAN),
    ("Poisson-disc 采样", "Bridson · 2007",
     "贪婪洗牌近似 O(n·k)\n避免扎堆的自然分布", PURPLE),
    ("L-System 晶体生长", "Prusinkiewicz & Lindenmayer · 1990",
     "球面坐标三维分支\n3-8 根侧生晶体 / 主晶", CYAN),
    ("Subsurface 散射", "Jensen et al. · 2001",
     "程序化点光源网络\n颜色匹配辉光 · 模拟 SSS", PURPLE),
]
cw = Inches(2.92); chh = Inches(4.6); gap = Inches(0.18)
sx = Inches(0.6); sy = Inches(1.5)
for i,(name,cite,desc,acc) in enumerate(cards):
    x = sx + i*(cw+gap)
    add_round(s, x, sy, cw, chh, fill=PANEL, line=acc, line_w=1.25, radius=0.06)
    # top accent strip
    add_rect(s, x, sy, cw, Inches(0.1), fill=acc)
    add_text(s, x+Inches(0.18), sy+Inches(0.28), cw-Inches(0.36), Inches(0.5),
             name, size=16, color=WHITE, bold=True)
    add_text(s, x+Inches(0.18), sy+Inches(0.82), cw-Inches(0.36), Inches(0.5),
             cite, size=10.5, color=acc, italic=True, font=FONT_EN)
    add_line(s, x+Inches(0.18), sy+Inches(1.3), cw-Inches(0.36), acc, 1.2)
    add_text(s, x+Inches(0.18), sy+Inches(1.5), cw-Inches(0.36), Inches(2.8),
             desc, size=13, color=GREY, line_spacing=1.35)
add_pagenum(s, 4)

# ---- SLIDE 5: 五阶段管线 ----
s = slide(5)
add_title(s, "五阶段全自动生成管线")
stages = [
    ("Stage 1", "地形生成", "Diamond-Square\n66K顶点 · 洞穴碗形"),
    ("Stage 2", "种子分布", "Poisson-disc\n面法线提取 · 间距保证"),
    ("Stage 3", "晶体生长", "六棱柱+尖端\nL-system 分支"),
    ("Stage 4", "灯光雾效", "三层布光\n体积雾 · 晶体 glow"),
    ("Stage 5", "相机动画", "960帧飞越\n生长动画 · DOF f/3.2"),
]
cw = Inches(2.15); chh = Inches(2.3); gap = Inches(0.32)
sx = Inches(0.6); sy = Inches(1.5)
for i,(st,name,desc) in enumerate(stages):
    x = sx + i*(cw+gap)
    acc = CYAN if i%2==0 else PURPLE
    add_round(s, x, sy, cw, chh, fill=PANEL, line=acc, line_w=1.2, radius=0.08)
    add_text(s, x, sy+Inches(0.12), cw, Inches(0.4), st, size=12, color=acc,
             bold=True, align=PP_ALIGN.CENTER, font=FONT_EN)
    add_text(s, x, sy+Inches(0.55), cw, Inches(0.5), name, size=17, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER)
    add_line(s, x+Inches(0.5), sy+Inches(1.15), cw-Inches(1.0), acc, 1.0)
    add_text(s, x+Inches(0.15), sy+Inches(1.3), cw-Inches(0.3), Inches(0.9),
             desc, size=11, color=GREY, align=PP_ALIGN.CENTER, line_spacing=1.3)
    # arrow between
    if i < 4:
        ax = x + cw + Inches(0.02)
        add_text(s, ax, sy+Inches(0.95), Inches(0.28), Inches(0.5), "→",
                 size=22, color=CYAN, bold=True, align=PP_ALIGN.CENTER,
                 font=FONT_EN)
# viewport screenshot at bottom
add_text(s, Inches(0.6), Inches(4.1), Inches(12), Inches(0.4),
         "◆ 管线输出 · Maya viewport 实景", size=13, color=CYAN, bold=True)
add_image_fit(s, "final_output/overview.0000.png",
              Inches(2.55), Inches(4.55), Inches(8.2), Inches(2.35))
add_pagenum(s, 5)

# ---- SLIDE 6: 代码架构 ----
s = slide(6)
add_title(s, "代码架构 — 7 个独立模块")
mods = [
    ("cave_terrain.py", "Diamond-Square 地形 · 洞穴围合 · 钟乳石/石笋 · 岩石材质"),
    ("crystal_seed.py", "Poisson-disc 采样 · 贪婪近似 O(n·k) · 面法线提取"),
    ("crystal_growth.py", "六棱柱+尖端 · L-system 球面分支 · aiNoise 云雾纹理"),
    ("cave_lighting.py", "三层布光 · aiAtmosphereVolume · 尘埃 · 动态相机"),
    ("crystal_ui.py", "Maya 内嵌 GUI · 4 标签页 16 参数 · 一键生成+渲染"),
    ("arnold_render.py", "AA 采样配置 · AOV · DOF 相机 · EXR 输出"),
    ("final_scene.py", "全管线集成 · 单函数 build() · 参数集中控制"),
]
add_text(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.45),
         "✦ 函数完整 docstring · 算法引用标注 · 单一职责 · 可独立测试 · 约 2,200 行",
         size=12.5, color=CYAN, italic=True)
cw = Inches(6.05); chh = Inches(0.66); gap = Inches(0.12)
sx1 = Inches(0.6); sx2 = Inches(6.85); sy = Inches(1.95)
for i,(name,desc) in enumerate(mods):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i-4
    x = sx1 if col==0 else sx2
    y = sy + row*(chh+gap)
    acc = CYAN if col==0 else PURPLE
    add_round(s, x, y, cw, chh, fill=PANEL, line=acc, line_w=0.9, radius=0.1)
    # icon chip
    add_rect(s, x+Inches(0.12), y+Inches(0.12), Inches(0.12), chh-Inches(0.24), fill=acc)
    add_text(s, x+Inches(0.35), y+Inches(0.04), Inches(2.7), chh, name,
             size=14, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE,
             font=FONT_MONO)
    add_text(s, x+Inches(3.0), y+Inches(0.04), cw-Inches(3.1), chh, desc,
             size=10.5, color=GREY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
add_pagenum(s, 6)

# ---- SLIDE 7: 晶体几何创新 ----
s = slide(7)
add_title(s, "晶体几何创新 — 从「钉子」到「水晶」")
cards = [
    ("⬡ 六棱柱", "polyCylinder(sx=6) 替代圆锥\n1.2–2.5m 半径 · 六方晶系", CYAN),
    ("↗ 法线对齐", "aimConstraint 精确对齐面法线\n替代手写欧拉角 · 避 gimbal lock", PURPLE),
    ("⎈ L-system 分支", "三维球面坐标方位角+仰角\n3–8 根侧生晶体 / 主晶", CYAN),
    ("⚓ 根部嵌入", "12% 嵌入消除悬浮\n底部 1.3× 径向扩展", PURPLE),
    ("◇ 粗细渐变", "主体粗 → 尖端细(0.92×)\npolyBevel 棱边倒角 · 10 种 IOR", CYAN),
]
cw = Inches(2.28); chh = Inches(2.95); gap = Inches(0.16)
sx = Inches(0.6); sy = Inches(1.55)
for i,(name,desc,acc) in enumerate(cards):
    x = sx + i*(cw+gap)
    add_round(s, x, sy, cw, chh, fill=PANEL, line=acc, line_w=1.1, radius=0.08)
    add_rect(s, x+Inches(0.0), sy, Inches(0.08), chh, fill=acc)
    add_text(s, x+Inches(0.22), sy+Inches(0.22), cw-Inches(0.4), Inches(0.6),
             name, size=16, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_line(s, x+Inches(0.22), sy+Inches(0.92), cw-Inches(0.5), acc, 1.0)
    add_text(s, x+Inches(0.22), sy+Inches(1.1), cw-Inches(0.44), Inches(1.7),
             desc, size=12, color=GREY, line_spacing=1.35)
# bottom note
add_round(s, Inches(0.6), Inches(4.75), Inches(12.15), Inches(0.85),
          fill=PANEL2, line=PURPLE_DIM, line_w=0.8, radius=0.1)
add_text(s, Inches(0.9), Inches(4.82), Inches(11.6), Inches(0.75),
         "◆ 创新：六棱柱六方晶系 + aimConstraint 法线对齐 + 球面 L-system 分支 — 共同构成「钉子」无法企及的真实晶簇形态。",
         size=12.5, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
add_pagenum(s, 7)

# ---- SLIDE 8: 材质与灯光 ----
s = slide(8)
add_title(s, "材质与灯光 — Arnold SSS + 体积雾")
# left: material bullets
lx = Inches(0.6)
add_text(s, lx, Inches(1.4), Inches(6.2), Inches(0.45),
         "✦ 10 种宝石 aiStandardSurface", size=15, color=CYAN, bold=True)
mat = [
    "Subsurface 次表面散射 (weight=0.8)",
    "Transmission 半透光玻璃 (0.6, depth=3.0)",
    "Coat 薄膜彩虹高光 (0.35, roughness=0.02)",
    "Dispersion 色散 (Abbe=45)",
    "Emission 晶体自发光 (0.35)",
    "10 种独立 IOR (1.43~1.77)",
    "aiNoise 3D 云雾纹理 → subsurfaceColor",
]
my = Inches(1.95)
for m in mat:
    add_text(s, lx+Inches(0.1), my, Inches(6.2), Inches(0.4),
             "▸ " + m, size=12.5, color=WHITE)
    my += Inches(0.43)
# right: lighting bullets
rx = Inches(7.1)
add_text(s, rx, Inches(1.4), Inches(5.7), Inches(0.45),
         "✦ 三层布光 + 体积雾 + 尘埃", size=15, color=PURPLE, bold=True)
lit = [
    "aiSkyDomeLight 暗紫环境光 (1.2)",
    "Key 暖金定向光 (120, -45°)",
    "Rim 冷蓝背光 (25, 135°)",
    "20 Crystal Glows (800–1600, 无衰减)",
    "aiAtmosphereVolume (density=0.008)",
    "120 暖色半透明尘埃 · DOF f/3.2",
    "ai_exposure=3.0 · 35mm 广角",
]
ly = Inches(1.95)
for l in lit:
    add_text(s, rx+Inches(0.1), ly, Inches(5.7), Inches(0.4),
             "▸ " + l, size=12.5, color=WHITE)
    ly += Inches(0.43)
# bottom render strip
add_round(s, Inches(0.6), Inches(5.05), Inches(12.15), Inches(1.55),
          fill=PANEL, line=CYAN_DIM, line_w=0.8, radius=0.06)
add_image_fit(s, "final_output/scene_latest/viewport.png",
              Inches(0.7), Inches(5.15), Inches(3.6), Inches(1.35))
add_text(s, Inches(4.5), Inches(5.3), Inches(8.1), Inches(1.2),
         "晶体自发光产生内部辉光 · 体积雾营造洞穴纵深 · 三层布光分离前景/背景/轮廓。\n视觉表现评估项：Arnold SSS + 体积雾 + 三层布光 — 真实矿物光学质感。",
         size=12.5, color=GREY, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)
add_pagenum(s, 8)

# ---- SLIDE 9: 交互 GUI ----
s = slide(9)
add_title(s, "交互 GUI — Maya 内嵌可视化窗口")
# mock window
wx = Inches(0.6); wy = Inches(1.45); ww = Inches(6.4); wh = Inches(5.2)
add_round(s, wx, wy, ww, wh, fill=RGBColor(0x16,0x13,0x32), line=CYAN,
          line_w=1.5, radius=0.03)
# title bar
add_rect(s, wx, wy, ww, Inches(0.42), fill=RGBColor(0x2A,0x24,0x55))
add_text(s, wx+Inches(0.2), wy+Inches(0.02), ww, Inches(0.4),
         "◆ Crystal Cavern v10", size=12, color=CYAN, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
# window dots
for i,c in enumerate([RGBColor(0xFF,0x5F,0x56),RGBColor(0xFF,0xBD,0x2E),RGBColor(0x27,0xC9,0x3F)]):
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, wx+ww-Inches(1.3)+i*Inches(0.22),
                           wy+Inches(0.13), Inches(0.16), Inches(0.16))
    d.fill.solid(); d.fill.fore_color.rgb = c; d.line.fill.background()
# tabs
tabs = ["Terrain","Crystals","Lighting","Render"]
tw = (ww-Inches(0.4))/4
for i,tn in enumerate(tabs):
    tx = wx+Inches(0.2)+i*tw
    active = (i==0)
    add_rect(s, tx, wy+Inches(0.5), tw-Inches(0.1), Inches(0.38),
             fill=PANEL2 if active else PANEL,
             line=CYAN if active else PANEL2, line_w=1.0)
    add_text(s, tx, wy+Inches(0.5), tw-Inches(0.1), Inches(0.38), tn,
             size=10.5, color=(CYAN if active else GREY), bold=active,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# controls
ctrls = ["Seed: 2026","Roughness: 0.45","Size: 40","Wall Height: 14",
         "Ceiling: 12","Density: 2.0","Variance: 1.0","Color Seed: 42",
         "Type Mix: Mixed","Ambient: 0.06","Fog: 0.008","Dust: 300",
         "Crystal Lights: 20","Resolution: 1920×1080","AA: 5"]
cy = wy+Inches(1.02)
for i,c in enumerate(ctrls):
    col = i%2
    row = i//2
    cx = wx+Inches(0.2)+col*( (ww-Inches(0.45))/2 )
    cyy = cy + row*Inches(0.27)
    add_rect(s, cx, cyy, (ww-Inches(0.55))/2, Inches(0.24),
             fill=RGBColor(0x12,0x0F,0x28), line=PANEL2, line_w=0.5)
    add_text(s, cx+Inches(0.08), cyy, (ww-Inches(0.55))/2, Inches(0.24),
             "▬ "+c, size=9.5, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
# buttons
bty = wy+wh-Inches(0.95)
add_round(s, wx+Inches(0.2), bty, Inches(2.0), Inches(0.4),
          fill=CYAN, line=None, radius=0.2)
add_text(s, wx+Inches(0.2), bty, Inches(2.0), Inches(0.4), "[ Generate ]",
         size=11, color=RGBColor(0x0f,0x0c,0x29), bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_round(s, wx+Inches(2.4), bty, Inches(1.6), Inches(0.4),
          fill=PANEL2, line=GREY_DIM, line_w=0.8, radius=0.2)
add_text(s, wx+Inches(2.4), bty, Inches(1.6), Inches(0.4), "[ Reset ]",
         size=10.5, color=GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_round(s, wx+Inches(4.2), bty, Inches(2.0), Inches(0.4),
          fill=PURPLE, line=None, radius=0.2)
add_text(s, wx+Inches(4.2), bty, Inches(2.0), Inches(0.4), "Render",
         size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, wx+Inches(0.2), wy+wh-Inches(0.45), ww-Inches(0.4), Inches(0.3),
         "progressWindow: 8 阶段进度条 (Cleaning → Materials → … → Lighting)",
         size=9.5, color=GREY_DIM, italic=True, align=PP_ALIGN.CENTER)
# right features
rx = Inches(7.3)
feats = [
    ("4 标签页", "Terrain/Crystals/Lighting/Render 分组", CYAN),
    ("16 参数", "全部可视化滑块 · 无需修改代码", PURPLE),
    ("一键生成", "30 秒完成全部 6 阶段管线", CYAN),
    ("一键渲染", "Arnold Render · 直接输出 EXR", PURPLE),
    ("可复现", "同 seed=相同洞穴 · 改参=新洞穴", CYAN),
    ("安全范围", "每控件 min/max · 防崩溃默认值", PURPLE),
]
fy = Inches(1.5)
for name,desc,acc in feats:
    add_round(s, rx, fy, Inches(5.45), Inches(0.78), fill=PANEL,
              line=acc, line_w=0.9, radius=0.1)
    add_text(s, rx+Inches(0.2), fy+Inches(0.06), Inches(2.0), Inches(0.32),
             name, size=13, color=acc, bold=True)
    add_text(s, rx+Inches(0.2), fy+Inches(0.4), Inches(5.1), Inches(0.32),
             desc, size=10.5, color=GREY)
    fy += Inches(0.88)
add_pagenum(s, 9)

# ---- SLIDE 10: 最终渲染 (scene1 vs scene2) ----
s = slide(10)
add_title(s, "最终渲染成果 — 两套参数对比")
add_text(s, Inches(0.6), Inches(1.4), Inches(6), Inches(0.45),
         "Scene 1 · 紫水晶洞穴", size=15, color=CYAN, bold=True)
add_image_fit(s, "final_output/scene1_amethyst/render.png",
              Inches(0.6), Inches(1.9), Inches(6.1), Inches(3.5))
add_text(s, Inches(0.6), Inches(5.5), Inches(6.1), Inches(0.4),
         "seed=2026 · 56 crystals · roughness=0.45", size=12, color=GREY,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(5.85), Inches(6.1), Inches(0.4),
         "9.7 MB EXR · 1920×1080 Arnold", size=10.5, color=GREY_DIM,
         align=PP_ALIGN.CENTER, italic=True)
add_text(s, Inches(6.65), Inches(1.4), Inches(6.1), Inches(0.45),
         "Scene 2 · 红宝石洞穴", size=15, color=PURPLE, bold=True)
add_image_fit(s, "final_output/scene2_ruby/render.png",
              Inches(6.65), Inches(1.9), Inches(6.1), Inches(3.5))
add_text(s, Inches(6.65), Inches(5.5), Inches(6.1), Inches(0.4),
         "seed=7777 · 56 crystals · roughness=0.55", size=12, color=GREY,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(6.65), Inches(5.85), Inches(6.1), Inches(0.4),
         "8.8 MB EXR · 1920×1080 Arnold", size=10.5, color=GREY_DIM,
         align=PP_ALIGN.CENTER, italic=True)
add_round(s, Inches(0.6), Inches(6.35), Inches(12.15), Inches(0.55),
          fill=PANEL2, line=CYAN_DIM, line_w=0.7, radius=0.1)
add_text(s, Inches(0.8), Inches(6.38), Inches(11.8), Inches(0.5),
         "◆ 每组参数生成完全不同的洞穴 — 同算法、不同 seed，形态/色彩/相机路径各异。",
         size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_pagenum(s, 10)

# ---- SLIDE 11: 动画成果 ----
s = slide(11)
add_title(s, "动画成果 — 飞越 + 生长")
# left card
add_round(s, Inches(0.6), Inches(1.5), Inches(6.05), Inches(5.0),
          fill=PANEL, line=CYAN, line_w=1.2, radius=0.05)
add_rect(s, Inches(0.6), Inches(1.5), Inches(6.05), Inches(0.1), fill=CYAN)
add_text(s, Inches(0.85), Inches(1.7), Inches(5.6), Inches(0.5),
         "960 帧飞越动画", size=20, color=CYAN, bold=True)
fly = [
    "✦ 动态相机路径 — 根据晶体位置自适应",
    "✦ 6 个关键帧 spline 插值 · 32 秒 (30fps)",
    "✦ 入口俯冲 → 紧贴晶体群掠过 → 全景上升",
    "✦ DOF 景深 f/3.2 · 35mm · 平滑样条",
    "✦ 每个 seed 生成不同飞行路径",
]
fy = Inches(2.45)
for f in fly:
    add_text(s, Inches(0.85), fy, Inches(5.6), Inches(0.45), f,
             size=13, color=WHITE, line_spacing=1.3)
    fy += Inches(0.62)
# big 960 number
add_text(s, Inches(0.85), Inches(5.55), Inches(5.6), Inches(0.8),
         "960 frames  ·  32s @ 30fps", size=18, color=PURPLE, bold=True,
         font=FONT_EN)
# right card
add_round(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(5.0),
          fill=PANEL, line=PURPLE, line_w=1.2, radius=0.05)
add_rect(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(0.1), fill=PURPLE)
add_text(s, Inches(7.1), Inches(1.7), Inches(5.4), Inches(0.5),
         "56 颗晶体生长动画", size=20, color=PURPLE, bold=True)
grow = [
    "✦ Scale 0→1 — 从裸岩长出成熟晶簇",
    "✦ 每颗错开启动帧 (spread 200 frames)",
    "✦ 生长周期 520 帧/颗 · 线性入 spline 出",
    "✦ 模拟天然矿物结晶过程",
    "✦ 0 帧=裸岩壁 · 720 帧=完整晶簇",
]
gy = Inches(2.45)
for g in grow:
    add_text(s, Inches(7.1), gy, Inches(5.4), Inches(0.45), g,
             size=13, color=WHITE, line_spacing=1.3)
    gy += Inches(0.62)
add_text(s, Inches(7.1), Inches(5.55), Inches(5.4), Inches(0.8),
         "56 crystals  ·  Scale 0 → 1", size=18, color=CYAN, bold=True,
         font=FONT_EN)
add_pagenum(s, 11)

# ---- SLIDE 12: AI工具 + 评估 + Thank you ----
s = slide(12)
add_title(s, "AI 工具使用 · 评估标准对照")
# left: AI tools
add_text(s, Inches(0.6), Inches(1.4), Inches(6), Inches(0.45),
         "✦ AI 工具使用记录", size=15, color=CYAN, bold=True)
tools = [
    ("Claude Code", "算法设计 · 代码生成 · 迭代优化 · 全模块重构"),
    ("Claude MCP", "Maya 直接操控 · 28 技能 · 场景构建 · 调试"),
    ("ChatGPT", "算法选型 · 论文引用核实 · API 查询"),
]
ty = Inches(1.95)
for name,desc in tools:
    add_round(s, Inches(0.6), ty, Inches(6.05), Inches(0.7), fill=PANEL,
              line=PURPLE_DIM, line_w=0.8, radius=0.1)
    add_text(s, Inches(0.8), ty+Inches(0.04), Inches(3.0), Inches(0.32),
             name, size=12.5, color=CYAN, bold=True, font=FONT_EN)
    add_text(s, Inches(0.8), ty+Inches(0.34), Inches(5.7), Inches(0.32),
             desc, size=10, color=GREY, line_spacing=1.2)
    ty += Inches(0.82)
# disclaimer
add_round(s, Inches(0.6), Inches(4.45), Inches(6.05), Inches(1.25),
          fill=PANEL2, line=PURPLE, line_w=0.7, radius=0.08)
add_text(s, Inches(0.8), Inches(4.55), Inches(5.7), Inches(1.1),
         "声明：AI 是实现工具而非创作主体。所有代码经 Maya 实机验证，\n设计决策（晶体形态、颜色方案、相机路径、参数范围）由作者做出。\n完整 AI 使用日志及提示词记录已保存。",
         size=10.5, color=GREY, line_spacing=1.35)
# right: assessment
add_text(s, Inches(6.95), Inches(1.4), Inches(5.8), Inches(0.45),
         "✦ 评估标准对照", size=15, color=PURPLE, bold=True)
assess = [
    ("算法复杂度", "4 个已发表论文 · 规范引用"),
    ("代码质量", "7 模块 · docstring · 注释完整"),
    ("GUI", "4 标签页 16 参数 · 一键生成+渲染"),
    ("视觉表现", "SSS + 体积雾 + 三层布光"),
    ("创新性", "矿物生成 · 六棱柱 · 云雾纹理"),
    ("AI 伦理", "完整使用日志 · 边界声明"),
]
ay = Inches(1.95)
for name,desc in assess:
    add_round(s, Inches(6.95), ay, Inches(5.8), Inches(0.62), fill=PANEL,
              line=CYAN_DIM, line_w=0.8, radius=0.1)
    add_text(s, Inches(7.15), ay+Inches(0.04), Inches(2.6), Inches(0.55),
             "✓ "+name, size=12.5, color=CYAN, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.9), ay+Inches(0.04), Inches(2.8), Inches(0.55),
             desc, size=10.5, color=GREY, anchor=MSO_ANCHOR.MIDDLE)
    ay += Inches(0.7)
# thank you footer
add_round(s, Inches(0.6), Inches(5.95), Inches(12.15), Inches(0.95),
          fill=PANEL2, line=CYAN, line_w=1.2, radius=0.06)
add_text(s, Inches(0.6), Inches(6.02), Inches(12.15), Inches(0.5),
         "Thank You", size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
         font=FONT_EN)
add_text(s, Inches(0.6), Inches(6.5), Inches(12.15), Inches(0.35),
         "Crystal Cavern  ·  未来创新设计 · 浙江大学 2026 · Prof. Xiaosong Yang",
         size=11, color=GREY, align=PP_ALIGN.CENTER)
add_pagenum(s, 12)

# ---------- save ----------
out = "答辩PPT_v3.pptx"
prs.save(out)
print("SAVED:", out)